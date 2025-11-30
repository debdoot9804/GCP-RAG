# backend/utils/parser.py

import tempfile
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation

def parse_file(filename: str, file_bytes: bytes) -> str:
    """Extract text from PDF, DOCX, PPTX, or TXT files (lightweight version)."""

    # Save file temporarily
    with tempfile.NamedTemporaryFile(delete=False, suffix=filename) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name

    text = ""

    # PDF
    if filename.lower().endswith(".pdf"):
        with fitz.open(tmp_path) as pdf:
            for page in pdf:
                text += page.get_text("text") + "\n"

    # DOCX
    elif filename.lower().endswith(".docx"):
        doc = Document(tmp_path)
        text = "\n".join([p.text for p in doc.paragraphs if p.text.strip()])

    # PPTX
    elif filename.lower().endswith(".pptx"):
        pres = Presentation(tmp_path)
        slides = []
        for slide in pres.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slides.append(shape.text)
        text = "\n".join(slides)

    # TXT
    elif filename.lower().endswith(".txt"):
        with open(tmp_path, "r", encoding="utf-8") as f:
            text = f.read()

    else:
        raise ValueError(f"❌ Unsupported file type: {filename}")

    if not text.strip():
        raise ValueError("⚠️ No text could be extracted from the file.")

    return text.strip()
